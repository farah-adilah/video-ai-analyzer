"""
Generation MCP Server
Handles PDF and PowerPoint report generation
"""
import asyncio
import json
from datetime import datetime
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.enums import TA_CENTER, TA_LEFT

class GenerationServer:
    def __init__(self):
        self.name = "generation-server"
        print(f"[{self.name}] Initialized")

    def _create_summary(self, analysis_data: dict) -> str:
        """
        Get AI-generated summary or fall back to simple summary
        """
        print(f"[{self.name}] Looking for summary in keys: {list(analysis_data.keys())}")
        
        if 'ai_summary' in analysis_data and analysis_data['ai_summary']:
            ai_summary = str(analysis_data['ai_summary'])
            
            skip_phrases = [
                "AI model not available",
                "Error",
                "Video analysis of",
                "completed successfully",
                "Analysis completed"
            ]
            
            if not any(phrase in ai_summary for phrase in skip_phrases):
                ai_summary = ai_summary.replace(
                    "Based on the video analysis data, the summary is:", ""
                ).strip()
                ai_summary = ai_summary.replace(
                    "Based on the video analysis data, a concise 2-3 sentence summary of the video content is as follows:", ""
                ).strip()
                
                lines = ai_summary.split('\n')
                clean_lines = []
                for line in lines:
                    line = line.strip()
                    skip_line = any(skip in line for skip in [
                        "Video duration:",
                        "Resolution:",
                        "Narration:",
                        "Visual scenes:"
                    ])
                    if not skip_line and line and len(line) > 10:
                        clean_lines.append(line)
                
                if clean_lines:
                    result = ' '.join(clean_lines)
                    print(f"[{self.name}] Using AI summary: {result[:100]}...")
                    return result
        
        print(f"[{self.name}] AI summary not found or invalid, using fallback")
        return self._create_simple_fallback(analysis_data)

    def _create_simple_fallback(self, analysis_data: dict) -> str:
        """Create simple summary when AI is not available"""
        parts = []
        
        if 'video_info' in analysis_data:
            duration = analysis_data['video_info'].get('duration', 0)
            parts.append(f"This {duration:.1f}-second video")
        
        if 'transcription' in analysis_data:
            trans = analysis_data['transcription']
            if isinstance(trans, dict) and trans.get('transcription'):
                trans_text = trans['transcription'][:80]
                parts.append(f"discusses: \"{trans_text}...\"")
            elif isinstance(trans, str) and trans:
                trans_text = trans[:80]
                parts.append(f"includes spoken content: \"{trans_text}...\"")
        
        if 'vision' in analysis_data:
            vision = analysis_data['vision']
            if vision.get('captions'):
                num_scenes = len(vision['captions'])
                captions = [c.get('caption', '') for c in vision['captions'][:2] if c.get('caption')]
                if captions:
                    scenes_desc = ' and '.join(captions)
                    parts.append(f"The video shows {num_scenes} key scenes including: {scenes_desc}")
        
        if parts:
            return '. '.join(parts) + '.'
        
    
        if 'video_name' in analysis_data:
            return f"Video analysis of {analysis_data['video_name']} completed successfully."
        
        return "Video content analysis completed."
    
    def _get_transcription_text(self, analysis_data: dict) -> str:
        """Get transcription text or appropriate message"""
        if 'transcription' not in analysis_data:
            return "No transcription available."
        
        trans = analysis_data['transcription']
        
        if isinstance(trans, dict):
            if trans.get('error'):
                return f"Transcription error: {trans['error']}"
            elif trans.get('transcription') and trans['transcription'].strip():
                return trans['transcription']
            else:
                return "No speech detected in this video. The video may be silent or contain only background noise/music."
        elif isinstance(trans, str) and trans.strip():
            return trans
        else:
            return "No speech detected in this video. The video may be silent or contain only background noise/music."
        
    async def generate_pptx(self, analysis_data: dict, output_path: str) -> dict:
        """
        Generate PowerPoint presentation
        
        Args:
            analysis_data: Analysis results
            output_path: Path to save PPTX
            
        Returns:
            Result dictionary
        """
        print(f"[{self.name}] Generating PowerPoint: {output_path}")
        
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            
            prs = Presentation()
            
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = title_slide.shapes.title
            subtitle = title_slide.placeholders[1]
            
            title.text = "Video Analysis Report"
            subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
            
            summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
            summary_slide.shapes.title.text = "Executive Summary"
            summary_text = self._create_summary(analysis_data)
            summary_slide.placeholders[1].text = summary_text

            if 'video_name' in analysis_data:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = "Video Information"
                content = slide.placeholders[1].text_frame
                content.text = f"Filename: {analysis_data['video_name']}"

            if 'transcription' in analysis_data:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = "Transcription"
                trans_text = self._get_transcription_text(analysis_data)
                slide.placeholders[1].text = trans_text
                
            prs.save(output_path)
            print(f"[{self.name}] PowerPoint saved: {output_path}")
            
            return {
                "success": True,
                "path": output_path,
                "message": "PowerPoint generated successfully"
            }
            
        except Exception as e:
            print(f"[{self.name}] Error generating PowerPoint: {e}")
            return {
                "success": False,
                "error": str(e)
            }
    
    async def generate_pdf(self, analysis_data: dict, output_path: str) -> dict:
        """Generate PDF report"""
        print(f"[{self.name}] Generating PDF: {output_path}")
        try:
            doc = SimpleDocTemplate(output_path, pagesize=letter)
            story = []
            styles = getSampleStyleSheet()
            
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=24,
                textColor='darkblue',
                spaceAfter=30,
                alignment=TA_CENTER
            )

            heading_style = ParagraphStyle(
                'CustomHeading',
                parent=styles['Heading2'],
                fontSize=16,
                textColor='darkgreen',
                spaceAfter=12,
                spaceBefore=12
            )
            
            story.append(Paragraph("Video Analysis Report", title_style))
            story.append(Spacer(1, 0.2*inch))
            story.append(Paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", styles['Normal']))
            story.append(Spacer(1, 0.5*inch))

            story.append(Paragraph("Executive Summary", heading_style))
            summary_text = self._create_summary(analysis_data)
            summary_text = summary_text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            story.append(Paragraph(summary_text, styles['BodyText']))
            story.append(Spacer(1, 0.5*inch))

            if 'video_name' in analysis_data:
                story.append(Paragraph("Video Information", heading_style))
                story.append(Paragraph(f"<b>Filename:</b> {analysis_data['video_name']}", styles['Normal']))
                
                if 'video_info' in analysis_data:
                    info = analysis_data['video_info']
                    story.append(Paragraph(f"<b>Resolution:</b> {info.get('width')}x{info.get('height')}", styles['Normal']))
                    story.append(Paragraph(f"<b>Duration:</b> {info.get('duration', 0):.1f} seconds", styles['Normal']))
                    story.append(Paragraph(f"<b>FPS:</b> {info.get('fps', 0):.1f}", styles['Normal']))
                
                story.append(Spacer(1, 0.3*inch))

            if 'transcription' in analysis_data:
                story.append(Paragraph("Transcription", heading_style))
               
                trans_text = self._get_transcription_text(analysis_data)
                trans_text = trans_text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                
                if "No speech detected" in trans_text or "error" in trans_text.lower():
                    story.append(Paragraph(f"<i>{trans_text}</i>", styles['Normal']))
                else:
                    story.append(Paragraph(trans_text, styles['Normal']))
                
                story.append(Spacer(1, 0.3*inch))

            doc.build(story)

            print(f"[{self.name}] PDF saved: {output_path}")
        
            return {
                "success": True,
                "path": output_path,
                "message": "PDF generated successfully"
            }

        except Exception as e:
            print(f"[{self.name}] Error generating PDF: {e}")
            import traceback
            traceback.print_exc()
            return {
                "success": False,
                "error": str(e)
            }

async def test():
    server = GenerationServer()
    test_data = {
        "video_name": "test.mp4",
        "transcription": "This is a test transcription",
        "captions": ["Scene 1", "Scene 2"]
    }
    result = await server.generate_pptx(test_data, "test_output.pptx")
    print(f"Result: {result}")

if __name__ == "__main__":
    print("Testing Generation Server...")
    asyncio.run(test())